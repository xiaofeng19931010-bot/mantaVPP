import sys
import os
# Add user site-packages to path if not present
user_site = os.path.expanduser("~/Library/Python/3.9/lib/python/site-packages")
if user_site not in sys.path:
    sys.path.append(user_site)

from pptx import Presentation

try:
    prs = Presentation(sys.argv[1])
    for i, slide in enumerate(prs.slides):
        print(f"\n--- Slide {i+1} ---")
        # Extract text from shapes
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(shape.text)
            # Also check for tables
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = [cell.text for cell in row.cells]
                    print(" | ".join(row_text))
except Exception as e:
    print(f"Error: {e}")
